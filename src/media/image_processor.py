import os
import re
import requests
import cv2
import numpy as np
from PIL import Image, ImageEnhance
import PyPDF2
from io import BytesIO
import base64
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from selenium.webdriver.common.action_chains import ActionChains
import pytesseract
from config import (POPPLER_PATH,TEMP_IMAGE_PATH,TCONNECT_HEADERS,TCONNECT_API_URL,TEMP_PDF_PATH, TEMP_WORD_PATH, TEMP_PPTX_PATH,tesseract_path)
from datetime import datetime
import time
from selenium.webdriver.common.by import By
from config import INTERFACE_KEYWORDS, EDUCATION_KEYWORDS
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.common.exceptions import NoSuchElementException, TimeoutException
from docx.oxml.ns import qn
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pdf2image.exceptions import PDFInfoNotInstalledError

pytesseract.pytesseract.tesseract_cmd = tesseract_path


def extract_text_from_image(image):
    try:
        image = image.convert('L')
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2.0)
        text = pytesseract.image_to_string(image, lang='eng', config='--psm 6').strip()
        print(f"OCR extracted: {text[:50]}...")
        return text
    except Exception as e:
        print(f"OCR error: {e}")
        return ""
 
def extract_text_from_pdf(pdf_path):
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + " "
                else:
                    print("Warning: Could not extract text from a PDF page using PyPDF2.")
            text = text.strip()
            if not text or "Dummy Information" in text:
                print("PyPDF2 extraction failed or returned dummy data, falling back to OCR.")
            else:
                print(f"PyPDF2 extracted: {text[:50]}...")
                return text
 
        images = convert_from_path(pdf_path, poppler_path=POPPLER_PATH)
        text = ""
        for i, image in enumerate(images):
            image = image.convert('L')
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(2.0)
            page_text = pytesseract.image_to_string(image, lang='eng', config='--psm 6').strip()
            if page_text:
                text += page_text + " "
            else:
                print(f"Warning: Could not extract text from PDF page {i+1} using OCR.")
        text = text.strip()
        if not text:
            print("Warning: OCR extracted empty text from PDF.")
            return "No text extracted from PDF"
        print(f"OCR extracted from PDF: {text[:50]}...")
        return text
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return "Failed to extract PDF text"
 
def extract_text_from_word(doc_path):
    try:
        if doc_path.endswith('.doc'):
            print("Warning: .doc file detected. Extraction not supported directly. Please convert to .docx or install 'antiword' for .doc support.")
            return "Failed to extract text from .doc file"
       
        doc = Document(doc_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + " "
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
        text = text.strip()
        if not text:
            print("Warning: No text extracted from Word document.")
            return "No text extracted from Word document"
        print(f"Word doc extracted: {text[:50]}...")
        return text
    except Exception as e:
        print(f"Word doc extraction error: {e}")
        return "Failed to extract text from Word document"
 
def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        text = text.strip()
        if not text:
            print("Warning: No text extracted from PPTX file.")
            return "No text extracted from PPTX file"
        print(f"PPTX extracted: {text[:50]}...")
        return text
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return "Failed to extract text from PPTX file"
 
def download_image(driver, img_element):
    try:
        img_src = img_element.get_attribute('src')
        if img_src.startswith('blob:'):
            image_blob_data = driver.execute_async_script("""
                const url = arguments[0];
                const callback = arguments[1];
                fetch(url)
                    .then(response => response.blob())
                    .then(blob => {
                        const reader = new FileReader();
                        reader.onloadend = () => callback(reader.result);
                        reader.readAsDataURL(blob);
                    });
            """, img_src)
            match = re.match(r'data:image/(?P<ext>.*?);base64,(?P<data>.*)', image_blob_data, re.DOTALL)
            if match:
                data = match.group('data')
                binary_data = base64.b64decode(data)
                img = Image.open(BytesIO(binary_data))
                img.save(TEMP_IMAGE_PATH)
                print(f"Image saved to {TEMP_IMAGE_PATH}")
                return img
        else:
            response = requests.get(img_src, timeout=5)
            img = Image.open(BytesIO(response.content))
            img.save(TEMP_IMAGE_PATH)
            print(f"Image saved to {TEMP_IMAGE_PATH}")
            return img
    except Exception as e:
        print(f"Error downloading image: {e}")
        return None
 
def download_document(driver, doc_element, message_content="", doc_type="pdf"):
    retries = 3
    original_window = driver.current_window_handle
    if doc_type == "pdf":
        default_filename = "unknown.pdf"
        temp_path = TEMP_PDF_PATH
    elif doc_type == "pptx":
        default_filename = "unknown.pptx"
        temp_path = TEMP_PPTX_PATH
    else:
        default_filename = "unknown.docx"
        temp_path = TEMP_WORD_PATH
    doc_file_name = default_filename
    download_dir = os.path.dirname(temp_path)
 
    # Extract filename from message content if available
    if message_content:
        pattern = r'([^\n]*\.pdf)' if doc_type == "pdf" else (r'([^\n]*\.pptx)' if doc_type == "pptx" else r'([^\n]*\.(doc|docx))')
        doc_name_match = re.search(pattern, message_content, re.IGNORECASE)
        if doc_name_match:
            doc_file_name = doc_name_match.group(1).strip()
            print(f"Extracted {doc_type.upper()} filename from message content: {doc_file_name}")
 
    downloaded_doc_path = os.path.join(download_dir, doc_file_name)
 
    for attempt in range(retries):
        try:
            # Hover over the message to ensure download button visibility
            ActionChains(driver).move_to_element(doc_element).perform()
            print(f"Hovered over message element for {doc_type.upper()} download.")
 
            # Updated XPaths to include div with role="button" and title containing "Download"
            possible_link_xpaths = [
                f".//div[@role='button' and contains(@title, 'Download')]",  # Matches log's div structure
                f".//a[contains(@href, '.{doc_type}') or @data-testid='media-url' or @data-testid='document-link']",
                ".//span[contains(@data-icon, 'doc')]/..",
                ".//div[contains(@class, 'media-content')]//a",
                f".//div[contains(text(), '.{doc_type}') or contains(@title, '.{doc_type}')]//a",
                f".//span[contains(text(), '.{doc_type}')]/ancestor::div[contains(@class, 'message-in')]//a",
                f".//div[contains(@class, 'message-in')]//div[contains(text(), '.{doc_type}') or contains(@title, 'document')]"
            ]
 
            doc_link = None
            for xpath in possible_link_xpaths:
                try:
                    doc_link = WebDriverWait(doc_element, 5).until(
                        EC.element_to_be_clickable((By.XPATH, xpath))
                    )
                    print(f"Found {doc_type.upper()} element with XPath: {xpath}")
                    break
                except (NoSuchElementException, TimeoutException):
                    continue
 
            if not doc_link:
                print(f"No clickable {doc_type.upper()} link found after checking all XPaths.")
                return False, doc_file_name
 
            # Log element details
            href = doc_link.get_attribute("href") or "No href"
            title = doc_link.get_attribute("title") or "No title"
            data_icon = doc_link.get_attribute("data-icon") or "No data-icon"
            print(f"{doc_type.upper()} link details - href: {href}, title: {title}, data-icon: {data_icon}")
 
            # Update filename from title or href if default is still used
            if title and title != "No title" and doc_file_name == default_filename:
                doc_file_name = title.replace("Download ", "").replace("&quot;", "").strip()
                print(f"Extracted {doc_type.upper()} filename from title: {doc_file_name}")
            elif href and href != "No href" and href.endswith(f'.{doc_type}') and doc_file_name == default_filename:
                doc_file_name = href.split('/')[-1]
                print(f"Extracted {doc_type.upper()} filename from href: {doc_file_name}")
 
            # Handle blob URLs
            blob_url = None
            if href.startswith("blob:"):
                blob_url = href
            else:
                blob_url = driver.execute_script("""
                    const link = arguments[0];
                    let blobUrl = null;
                    if (link.href && link.href.startsWith('blob:')) {
                        blobUrl = link.href;
                    } else {
                        const media = link.closest('div[data-testid="media-content"]');
                        if (media) {
                            const source = media.querySelector('source');
                            if (source && source.src && source.src.startsWith('blob:')) {
                                blobUrl = source.src;
                            }
                        }
                    }
                    return blobUrl;
                """, doc_link)
 
            if blob_url:
                mime_type = (
                    "application/pdf" if doc_type == "pdf" else
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation" if doc_type == "pptx" else
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
                doc_data = driver.execute_async_script("""
                    const url = arguments[0];
                    const callback = arguments[1];
                    fetch(url)
                        .then(response => response.blob())
                        .then(blob => {
                            const reader = new FileReader();
                            reader.onloadend = () => callback(reader.result);
                            reader.readAsDataURL(blob);
                        })
                        .catch(err => callback(null));
                """, blob_url)
               
                if doc_data:
                    match = re.match(rf'data:{mime_type};base64,(?P<data>.*)', doc_data, re.DOTALL)
                    if match:
                        data = match.group('data')
                        binary_data = base64.b64decode(data)
                        with open(temp_path, 'wb') as f:
                            f.write(binary_data)
                        print(f"{doc_type.upper()} successfully downloaded to {temp_path} via blob URL")
                        return True, doc_file_name
                    else:
                        print(f"Failed to parse {doc_type.upper()} blob data")
                else:
                    print(f"Failed to fetch {doc_type.upper()} blob URL: {blob_url}")
 
            # Try clicking the link using ActionChains
            ActionChains(driver).move_to_element(doc_link).click().perform()
            print(f"Clicked {doc_type.upper()} link on attempt {attempt + 1} using ActionChains")
 
            # Wait for file to download
            wait_time = 0
            max_wait = 20  # Increased wait time for live environment
            while wait_time < max_wait:
                if os.path.exists(downloaded_doc_path) and os.path.getsize(downloaded_doc_path) > 0:
                    os.rename(downloaded_doc_path, temp_path)
                    print(f"{doc_type.upper()} successfully downloaded to {temp_path}")
                    return True, doc_file_name
                for file in os.listdir(download_dir):
                    if file.lower().endswith(f".{doc_type}"):
                        found_path = os.path.join(download_dir, file)
                        if os.path.getsize(found_path) > 0:
                            os.rename(found_path, temp_path)
                            print(f"{doc_type.upper()} found with name {file}, renamed to {temp_path}")
                            return True, doc_file_name
                time.sleep(1)
                wait_time += 1
            print(f"{doc_type.upper()} file not found at {temp_path} or {downloaded_doc_path} after {max_wait} seconds")
 
        except Exception as e:
            print(f"Error downloading {doc_type.upper()} on attempt {attempt + 1}: {e}")
            if len(driver.window_handles) > 1:
                driver.close()
                driver.switch_to.window(original_window)
            time.sleep(2)
 
    print(f"Failed to download {doc_type.upper()} after {retries} attempts")
    return False, doc_file_name

def create_tconnect_task(contact_id, interface):
    """
    Create a t-connect task for CMS or LMS issues and return the formatted task ID.
    
    Args:
        contact_id (str): The contact ID or application number of the student.
        interface (str): The interface (CMS or LMS).
    
    Returns:
        tuple: (success: bool, task_id: str or None) - Formatted task ID (e.g., '#116157') if successful, None otherwise.
    """
    if interface not in ["CMS", "LMS"]:
        print(f"❌ Invalid interface: {interface}. Only CMS and LMS are supported.")
        return False, None
    
    description = (
        f"Student Contact: {contact_id}\n"
        f"Application: {interface}\n"
        f"Issue: The student is experiencing difficulties logging into the portal. "
        f"Immediate assistance is requested to resolve this access issue.\n"
        f"Kindly solve the issue at the earliest convenience.\n\n"
        f"Thank you\nDataTeam(AI bot)"
    )
    
    payload = {
        "category": "Student Grievance",
        "description": description,
        "task_date": datetime.now().strftime("%Y-%m-%d"),
        "tat": "1",
        "responsible_email": "sownder.a@texila.org",
        "created_email": "AISMSupport@texila.org",
        "request_from": "data_team"
    }
    
    retries = 3
    for attempt in range(retries):
        try:
            response = requests.post(TCONNECT_API_URL, json=payload, headers=TCONNECT_HEADERS, timeout=15)
            print(f"Attempt {attempt + 1}: API response status: {response.status_code}, Response: {response.text[:100]}...")
            if response.status_code == 200:
                try:
                    response_data = response.json()
                    task_id = response_data.get("tconnect_task_id")
                    if task_id:
                        formatted_task_id = f"#{task_id}"
                        print(f"✅ Task created for {contact_id} with Task ID: {formatted_task_id}")
                        return True, formatted_task_id
                    else:
                        print(f"❌ Task created but no tconnect_task_id found in response: {response.text}")
                        if attempt < retries - 1:
                            print(f"Retrying task creation (attempt {attempt + 2}/{retries})...")
                            time.sleep(2 ** attempt)  # Exponential backoff: 1s, 2s, 4s
                        continue
                except ValueError as json_err:
                    print(f"❌ Invalid JSON response: {response.text[:100]}... Error: {json_err}")
                    if attempt < retries - 1:
                        print(f"Retrying task creation (attempt {attempt + 2}/{retries})...")
                        time.sleep(2 ** attempt)
                    continue
            else:
                print(f"❌ Failed to create task: Status {response.status_code}, Response: {response.text[:100]}...")
                if attempt < retries - 1:
                    print(f"Retrying task creation (attempt {attempt + 2}/{retries})...")
                    time.sleep(2 ** attempt)
                continue
        except requests.exceptions.RequestException as e:
            print(f"❌ API error on attempt {attempt + 1}: {e}")
            if attempt < retries - 1:
                print(f"Retrying task creation (attempt {attempt + 2}/{retries})...")
                time.sleep(2 ** attempt)
    
    print(f"❌ Failed to create task for {contact_id} after {retries} attempts")
    return False, None

def is_unsaved_contact(contact_name, contact_number):
    """
    Determine if a contact is unsaved by checking if the contact name or number starts with '+'.
    
    Args:
        contact_name (str): The contact name as displayed in WhatsApp.
        contact_number (str): The extracted contact number.
    
    Returns:
        bool: True if the contact is unsaved (starts with '+'), False otherwise.
    """
    # Check if contact_name or contact_number starts with '+'
    if contact_name.startswith('+') or contact_number.startswith('+'):
        return True
    # Additional check for saved contacts with application number format (e.g., TGY/12345, TZM/12345, TOC/12345, or 12345)
    app_number_pattern = r'^(?:[A-Z]{3}/)?\d{5}$'
    if re.match(app_number_pattern, contact_name) or re.match(app_number_pattern, contact_number):
        return False
    return False  # Default to saved contact if no clear unsaved indicator

def is_general_media(text, media_type="image"):
    """
    Determine if a media (image, PDF, Word, PPTX) is general by checking if its extracted text
    lacks LMS, CMS, remittance, or education-related keywords.
    
    Args:
        text (str): The extracted text from the media.
        media_type (str): Type of media ('image', 'pdf', 'word', 'pptx').
    
    Returns:
        bool: True if the media is general (no specific keywords or extraction failure), False otherwise.
    """
    if not text or any(error in text.lower() for error in ["failed to extract", "no text extracted", "lms cms ui issue", "processing issue"]):
        return True  # Treat empty text or extraction errors as general
    text_lower = text.lower()
    # Define keywords for CMS, LMS, remittance, and education
    specific_keywords = (
        EDUCATION_KEYWORDS +
        INTERFACE_KEYWORDS['lms'] +
        INTERFACE_KEYWORDS['cms'] +
        INTERFACE_KEYWORDS['remittance'] +
        ["learning management", "course management", "portal", "login", "dashboard", "module", "assignment",
         "payment", "invoice", "transaction", "bank", "receipt", "finance", "fee", "due"]
    )
    return not any(keyword in text_lower for keyword in specific_keywords)

def detect_interface(text):
    text_lower = text.lower()
    cms_keywords = INTERFACE_KEYWORDS['cms']
    lms_keywords = INTERFACE_KEYWORDS['lms']
    cms_count = sum(1 for kw in cms_keywords if kw in text_lower)
    lms_count = sum(1 for kw in lms_keywords if kw in text_lower)
    if cms_count >= 2:
        return "CMS"
    elif lms_count >= 2:
        return "LMS"
    return None

